import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE

# Color Constants
COLOR_BG = RGBColor(10, 13, 20)           # Deep Space Black/Blue
COLOR_CARD_BG = RGBColor(19, 25, 36)      # Sleek Slate Blue Card Background
COLOR_CARD_BORDER = RGBColor(30, 41, 59)  # Slate Gray Border
COLOR_NEON_ACCENT = RGBColor(0, 240, 255) # Cyber Cyan Accent
COLOR_TITLE = RGBColor(255, 255, 255)     # Pure White for titles and headers
COLOR_BODY = RGBColor(168, 178, 209)      # Cool Silver-Gray for body descriptions

def make_backup(pptx_path):
    dir_name = os.path.dirname(pptx_path)
    base_name = os.path.basename(pptx_path)
    backup_path = os.path.join(dir_name, "Predictra-Threat-Map-Backup.pptx")
    shutil.copy2(pptx_path, backup_path)
    print(f"Backup created at: {backup_path}")

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_transition(slide):
    # Transition XML for a slow fade transition
    transition_xml = '''
    <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow">
        <p:fade />
    </p:transition>
    '''
    slide.element.insert(-1, parse_xml(transition_xml))

def add_card_to_back(slide, left, top, width, height):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_CARD_BORDER
    card.line.width = Pt(1)
    
    # Send shape to back of the slide hierarchy so text overlays correctly
    slide.shapes._spTree.remove(card._element)
    slide.shapes._spTree.insert(2, card._element)
    return card

def style_general_shape(shape, is_title=False):
    # Apply style to card background shapes if applicable
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # If it's a solid card background shape (width and height are substantial)
            if shape.width > Inches(1.5) and shape.height > Inches(0.5) and not is_title:
                shape.fill.solid()
                shape.fill.fore_color.rgb = COLOR_CARD_BG
                shape.line.color.rgb = COLOR_CARD_BORDER
                shape.line.width = Pt(1)
            else:
                # Transparent for headers and title shapes
                shape.fill.background()
                shape.line.fill.background()
    except:
        pass

    if shape.has_text_frame:
        shape.text_frame.margin_left = Inches(0.1)
        shape.text_frame.margin_right = Inches(0.1)
        shape.text_frame.margin_top = Inches(0.1)
        shape.text_frame.margin_bottom = Inches(0.1)
        
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            
            is_chapter = "CHAPTER" in text.upper()
            
            if not paragraph.runs:
                paragraph.add_run()
                
            for run in paragraph.runs:
                run.font.name = "Segoe UI"
                if is_title:
                    run.font.size = Pt(28)
                    run.font.bold = True
                    run.font.color.rgb = COLOR_TITLE
                elif is_chapter:
                    run.font.size = Pt(10)
                    run.font.bold = True
                    run.font.color.rgb = COLOR_NEON_ACCENT
                elif len(text) <= 35 and '\n' not in text:
                    # Sub-headers
                    run.font.size = Pt(13)
                    run.font.bold = True
                    run.font.color.rgb = COLOR_TITLE
                else:
                    # Body descriptions
                    run.font.size = Pt(11)
                    run.font.bold = False
                    run.font.color.rgb = COLOR_BODY

def enhance_presentation(pptx_path, new_usecase_path, new_class_path):
    prs = Presentation(pptx_path)
    
    # Speaker Notes Dictionary (0-indexed for 26 slides)
    SPEAKER_NOTES = {
        0: ("Brahim: Good morning, ladies and gentlemen of the jury. Welcome to our graduation defense for our "
            "Bachelor’s Degree in Information Technology. My name is Brahim Jaballi, and alongside my colleague "
            "Chiheb Amri, we are proud to present our final year project: the 'Predictra Threat Map'. This project, "
            "conducted under the supervision of Mr. Anis Dhahri at ISET Gafsa, and in collaboration with host company "
            "Predictra Cybersecurity, focuses on the design and development of a real-time spatialized threat "
            "intelligence visualization platform. We would also like to thank Ms. Hadhami Issaoui and Mr. Walid Hammemi "
            "for presiding and reporting on our defense today."),
        
        1: ("Brahim: To guide you through our presentation, we have structured our defense into five key chapters. "
            "I will begin with the Introduction and the Pre-Study, outlining the cybersecurity landscape, our problematic, "
            "and the Agile Scrum methodology. Next, we will present our Conceptual Study, including system actors, "
            "requirements, use cases, and architecture. We will then dive into the Realization phase, detailing "
            "Sprints 1 to 4 and showcasing the system interfaces. Finally, we will conclude with a retrospective, "
            "technical challenges resolved, and future perspectives."),
            
        2: ("Brahim: Let us start with the introduction. Today, corporate networks and cloud platforms are faced "
            "with an expanded attack surface. Highly coordinated threat groups launch disruptive operations globally. "
            "As Sun Tzu wrote in The Art of War: 'If you know the enemy and know yourself, you need not fear the "
            "result of a hundred battles.' In cybersecurity, this means shifting from a reactive posture to a "
            "proactive threat intelligence model, where we understand external adversary infrastructure before a breach occurs."),
            
        3: ("Brahim: Our host organization, Predictra Cybersecurity, is an AI-driven Cyber Threat Intelligence "
            "firm. Their goal is to democratize threat intelligence, transforming raw feeds into actionable visual "
            "insights. A critical industry metric is the compromise dwell time, which globally averages over 200 days. "
            "Security teams often detect breaches months after they occur. Our Threat Map aims to shrink this window "
            "from months to seconds by providing immediate, real-time visibility."),
            
        4: ("Brahim: The core problem we address is threefold. First, spreadsheet overload. SOC analysts suffer "
            "from visual fatigue reading thousands of static logs. Second, a lack of spatial context. Traditional "
            "dashboards fail to show attack origins, targets, and trajectories geographically. Third, data heterogeneity. "
            "Threat feeds are scattered across multiple vendors in incompatible formats. This makes manual correlation "
            "slow, leading to delayed responses."),
            
        5: ("Brahim: To solve this, we propose the Predictra Threat Map, built on a four-layer architecture. "
            "First, the Ingestion Layer harvests data asynchronously from 9+ feeds. Second, the Enrichment Layer "
            "geolocates attacks and classifies victims into 9 critical sectors. Third, the Streaming Layer broadcasts "
            "events using Server-Sent Events with a 300ms batching queue. Fourth, the Visualization Layer displays "
            "curved attack arcs on an interactive 3D WebGL Earth, alongside a STIX parser and MITRE ATT&CK grid."),
            
        6: ("Brahim: We adopted the Agile Scrum methodology, organizing our work into 4 sprints over a 12-week period. "
            "Sprint 1 focused on backend ingestion and schemas. Sprint 2 implemented the SSE stream and the 3D globe. "
            "Sprint 3 integrated the STIX 2.1 parser and MITRE ATT&CK grid. Sprint 4 added the history browser, "
            "Target My IP, and performance guardrails. We held twice-weekly rituals on Thursdays for code merges "
            "and profiling, and Sundays for architectural consolidation and documentation."),
            
        7: ("Brahim: This is our Global Use Case diagram, showing the system boundaries and interactions. The primary "
            "system actor is the set of External Threat Feeds, which automatically trigger backend ingestion. The "
            "human actor is the Security Analyst, who monitors the 3D globe and dashboard, filters threat logs, searches "
            "historical data, uploads STIX bundles, and exports reports. This diagram establishes the foundational "
            "boundaries of the Predictra Threat Map."),
            
        8: ("Brahim: To elaborate on the actors: the External Threat Feeds run continuously without human "
            "intervention, pushing Indicators of Compromise via HTTP REST, SSE, and WebSockets. On the other side, "
            "the Security Analyst interacts with the frontend dashboard. They can observe live threats, filter events "
            "by country or sector, analyze STIX files locally for privacy, and map IP details. This separation ensures "
            "that the core ingestion runs independently of analyst sessions."),
            
        9: ("Brahim: We defined 9 functional requirements, denoted RF-1 to RF-9. These include multi-feed ingestion, "
            "geolocation enrichment, target sector classification, real-time SSE streaming, the 3D WebGL globe, "
            "client-side STIX parsing, MITRE ATT&CK mapping, historical log searching, and the 'Target My IP' feature. "
            "Each requirement maps directly to a user story delivered in our Scrum backlog."),
            
        10: ("Brahim: Our non-functional requirements ensure the quality and stability of the platform. RNF-1 "
             "guarantees a 60fps fluidity on the 3D globe. RNF-2 uses a 10,000-event circular RingBuffer to cap memory "
             "usage. RNF-3 sets a 30-day TTL index in MongoDB to prune old logs. RNF-4 monitors database quota, "
             "switching to in-memory streaming if DB size exceeds 500MB. RNF-5 implements a modern, dark glassmorphism "
             "user interface to reduce visual strain."),
             
        11: ("Chiheb: Thank you, Brahim. I will now present our technology stack. For the backend runtime, we selected "
             "Node.js and Express 5 over Java Spring Boot due to its non-blocking asynchronous event loop, which is perfect "
             "for streaming SSE connections. For state management, we chose Zustand 5 over Redux for high-speed reactive "
             "updates outside the React DOM. For graphics, we used Three.js and React Three Fiber to build custom shaders. "
             "We chose Vite 7 for native ESM builds, and MongoDB Atlas for its schemaless BSON structure."),
             
        12: ("Chiheb: Our system architecture is split into three main tiers. The Frontend is built in React and "
             "served via Vercel CDN, utilizing API routing proxies to eliminate cross-origin resource sharing blocks. "
             "The Backend runs on a Cloud VM using PM2 process clustering across CPU cores for automatic crash recovery. "
             "The Database tier uses MongoDB Atlas with a replicated cluster, automatic TTL pruning, and size-based quota guards."),
             
        13: ("Chiheb: This is our Global Class Diagram, showing the structural design of the platform. In the Backend "
             "tier, the Express Server schedules multiple Scraper Services, which feed into the Enrichment Service to "
             "resolve geolocations and sectors. The Zustand Store manages frontend state with a circular RingBuffer and "
             "performance telemetry. The React Views, such as the Dashboard, STIX Workspace, and Analytics page, "
             "subscribe directly to Zustand slices for reactivity."),
             
        14: ("Brahim: In Sprint 1, we built the core ingestion pipeline. We integrated 9 feeds: Checkpoint, "
             "Bitdefender, Fortinet, URLhaus, Kaspersky Feodo, AlienVault, RansomWatch, MontySecurity, and MISP Galaxy. "
             "We designed a compact BSON schema in MongoDB, reducing the database footprint from 200 bytes to 80 bytes "
             "per document. This ensures long-term database efficiency and lower storage costs."),
             
        15: ("Brahim: During Sprint 1, we modeled the enrichment logic. The Enrichment Service uses a 5-layer "
             "classification engine to identify the targeted industry sector: (1) victim keywords, (2) event details "
             "regex, (3) destination port mapping, (4) MISP Galaxy metadata, and (5) scraper-specific defaults. We also "
             "integrated an asynchronous RDAP cache for IP owner lookups with a 3-second timeout to prevent API blocking."),
             
        16: ("Brahim: Sprint 2 focused on streaming and 3D math. We built the SSE broadcast engine using a 300ms "
             "flush scheduler that batches queued events into single payloads to prevent browser flooding. For the 3D "
             "globe, geographic coordinates are converted to Cartesian vectors. The attack arcs use Spherical Linear "
             "Interpolation, or SLERP, to curve around the Earth, with Bezier elevation lifting longer distance arcs higher."),
             
        17: ("Brahim: To make the globe visually stunning, we wrote custom GLSL shaders. The Atmospheric Glow "
             "uses a Fresnel-effect vertex and fragment shader. Neon Country Borders are rendered using 5 overlapping "
             "TopoJSON vector layers with decreasing opacities. Finally, Bloom Post-Processing applies a luminosity "
             "threshold to make severity-colored attack arcs glow like neon light fibers."),
             
        18: ("Chiheb: In Sprint 3, we built the STIX 2.1 Threat Workspace. The client-side JSON parser processes "
             "threat bundles entirely in the browser, ensuring sensitive files never leave the analyst's machine. "
             "The parser extracts STIX Domain Objects like Threat Actors, Campaigns, and Malware. These are automatically "
             "mapped to the 12 phases of the MITRE ATT&CK Kill Chain, updating the metrics instantly."),
             
        19: ("Chiheb: For threat visualization in Sprint 3, we implemented an interactive D3 Force Graph. It "
             "renders STIX objects as nodes and their relationships as edges. The simulation applies repulsion, "
             "spring, and gravity forces using Verlet integration until it cools into a stable layout. Analysts can hover "
             "over nodes to inspect raw SDO metadata and explore relationship chains."),
             
        20: ("Chiheb: Sprint 4 focused on UI dashboards, history, and guardrails. We implemented a paginated History "
             "Browser query with compound indexes executing under 50ms. 'Target My IP' allows analysts to geolocate "
             "themselves and auto-filter threat logs targeting their subnet. Excel Export generates client-side worksheets "
             "for up to 5,000 logs. Most importantly, we implemented performance guardrails to preserve frame rate."),
             
        21: ("Chiheb: Our performance engine uses an Adaptive Event Sampling loop. If the frame rate is high "
             "(above 55fps), all events are rendered. If load increases and FPS drops between 30 and 55, the drop rate "
             "is set to 50%. Under extreme stress (below 30fps), the engine drops 75% of visual events. Supported by "
             "the O(1) RingBuffer, this prevents memory leaks and garbage collection spikes."),
             
        22: ("Brahim: Here is an overview of our application interfaces. The 3D Live Threat Map displays the "
             "rotating Earth with glowing arcs. The System Dashboard provides real-time metrics, country statistics, "
             "and a 2D flat map. The STIX Intelligence Workspace features the D3 graph and the MITRE ATT&CK grid. "
             "Finally, the Analytics Dashboard aggregates threat metrics, sector distributions, and trends."),
             
        23: ("Brahim: Looking back, our team maintained a stable Scrum velocity of 25–30 story points per "
             "sprint, delivering all features. We resolved three critical challenges: visual lag via event batching "
             "and adaptive sampling; database limits via BSON compression and TTL indexes; and deployment blocks via "
             "Vercel reverse-proxy routing rules to eliminate CORS issues."),
             
        24: ("Chiheb: For future work, we have designed a three-step scalability roadmap. First, Kafka Distributed "
             "Ingestion to decouple scrapers from the server. Second, Predictive ML Extensions to forecast attack "
             "spikes and project heatmaps onto the globe. Third, Sharded Geo-Database Clustering to partition MongoDB "
             "by geographic region, minimizing latency and scaling horizontally."),
             
        25: ("Chiheb: In conclusion, the Predictra Threat Map successfully bridges the visibility gap in cyber "
             "defense. It integrates 9+ feeds, feeds a 3D WebGL globe at 60fps, and features a STIX 2.1 parser. It "
             "reduces dwell time from months to seconds. This project demonstrates a rigorous academic and engineering "
             "contribution. Thank you for your time, and we welcome any questions you may have.")
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Processing Slide {slide_idx+1}...")
        
        # Set Dark Background
        set_slide_background(slide)
        
        # Add Slide Fade Transition
        add_transition(slide)
        
        # Add Speaker Notes
        if slide_idx in SPEAKER_NOTES:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = SPEAKER_NOTES[slide_idx]
            
        # Determine Title Shape
        title_shape = slide.shapes.title
        
        # SLIDE-SPECIFIC CUSTOM EDITS
        # ===========================
        
        # Slide 8: Global Use Case Diagram
        if slide_idx == 7:
            # Center the new Use Case Diagram (width=2.55", height=5.0" based on aspect ratio 0.51)
            # Center coordinates on 10.0" x 5.625" slide: left=3.72", top=0.31"
            for shape in list(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide.shapes._spTree.remove(shape._element)
            slide.shapes.add_picture(new_usecase_path, Inches(3.72), Inches(0.31), Inches(2.55), Inches(5.0))
            
        # Slide 11: Non-Functional Requirements
        elif slide_idx == 10:
            # No diagrams here. We adjust title text
            if title_shape:
                title_shape.text_frame.text = "Non-Functional Requirements & Global Use Case"
            # Add a small note pointing to the usecase slide on Slide 8
            note_shape = slide.shapes.add_textbox(Inches(5.14), Inches(1.36), Inches(4.5), Inches(3.5))
            tf = note_shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Global Use Case Model Reference"
            p.font.name = "Segoe UI"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_TITLE
            
            p2 = tf.add_paragraph()
            p2.text = "\nTo view the complete system boundaries and user/system actor interactions, please refer to Slide 8.\n\nThe Use Case model illustrates:\n• Ingestion triggers by external thread feeds.\n• Investigation actions by SOC analysts (Filters, STIX Parsing, exports)."
            p2.font.name = "Segoe UI"
            p2.font.size = Pt(11)
            p2.font.color.rgb = COLOR_BODY
            
        # Slide 13: System Architecture
        elif slide_idx == 12:
            # 1. Remove all diagram images
            for shape in list(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide.shapes._spTree.remove(shape._element)
            
            # 2. Re-layout remaining text cards to occupy the full slide in three columns
            # Column 1 (Frontend): left=1.16, top=1.5
            # Column 2 (Backend): left=3.8, top=1.5
            # Column 3 (Database): left=6.44, top=1.5
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    try:
                        # Move background shapes based on original left positions
                        # Frontend Background (left around 1.16)
                        if abs(shape.left/914400 - 1.16) < 0.1:
                            shape.left = int(Inches(1.16))
                            shape.top = int(Inches(1.5))
                            shape.width = int(Inches(2.4))
                            shape.height = int(Inches(3.3))
                        # Backend Background (left around 3.72)
                        elif abs(shape.left/914400 - 3.72) < 0.1:
                            shape.left = int(Inches(3.8))
                            shape.top = int(Inches(1.5))
                            shape.width = int(Inches(2.4))
                            shape.height = int(Inches(3.3))
                        # Database Background (left around 6.28)
                        elif abs(shape.left/914400 - 6.28) < 0.1:
                            shape.left = int(Inches(6.44))
                            shape.top = int(Inches(1.5))
                            shape.width = int(Inches(2.4))
                            shape.height = int(Inches(3.3))
                    except:
                        pass
                else:
                    text = shape.text_frame.text
                    # Move text shapes by checking their content
                    if "Frontend (Vercel CDN)" in text:
                        shape.left = int(Inches(1.26))
                        shape.top = int(Inches(1.7))
                        shape.width = int(Inches(2.2))
                    elif "React SPA hosted globally" in text:
                        shape.left = int(Inches(1.26))
                        shape.top = int(Inches(2.1))
                        shape.width = int(Inches(2.2))
                        shape.height = int(Inches(2.5))
                    elif "Backend (Cloud VM / PM2)" in text:
                        shape.left = int(Inches(3.9))
                        shape.top = int(Inches(1.7))
                        shape.width = int(Inches(2.2))
                    elif "Node.js Express server" in text:
                        shape.left = int(Inches(3.9))
                        shape.top = int(Inches(2.1))
                        shape.width = int(Inches(2.2))
                        shape.height = int(Inches(2.5))
                    elif "Database (MongoDB Atlas)" in text:
                        shape.left = int(Inches(6.54))
                        shape.top = int(Inches(1.7))
                        shape.width = int(Inches(2.2))
                    elif "Replicated BSON cluster" in text:
                        shape.left = int(Inches(6.54))
                        shape.top = int(Inches(2.1))
                        shape.width = int(Inches(2.2))
                        shape.height = int(Inches(2.5))
            
        # Slide 14: Global Class Diagram
        elif slide_idx == 13:
            # Replace Class Diagram at same position
            for shape in list(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide.shapes._spTree.remove(shape._element)
            slide.shapes.add_picture(new_class_path, Inches(1.99), Inches(1.01), Inches(5.67), Inches(2.54))
            
        # Slide 16: Sprint 1 Diagrams
        elif slide_idx == 15:
            # 1. Remove picture shapes
            for shape in list(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide.shapes._spTree.remove(shape._element)
            
            # 2. Re-layout description to fill slide nicely and change title
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if "Sprint 1: Class & Sequence Diagrams" in text:
                        shape.text_frame.text = "Sprint 1: Enrichment Engine & Ingestion Logic"
                    elif "The EnrichmentService applies" in text:
                        shape.left = int(Inches(1.0))
                        shape.top = int(Inches(1.8))
                        shape.width = int(Inches(8.0))
                        shape.height = int(Inches(3.2))
                        # Increase body font size slightly since it is now a summary slide
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(14)
                                
        # Slide 18: Sprint 2 Shaders & Sequence
        elif slide_idx == 17:
            # 1. Remove sequence diagram empty placeholders (top <= 1.5, left <= 5.0)
            for shape in list(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.left/914400 < 5.0 and shape.top/914400 > 1.2 and shape.top/914400 < 1.6:
                    slide.shapes._spTree.remove(shape._element)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide.shapes._spTree.remove(shape._element)
            
            # 2. Update Slide Title
            for shape in slide.shapes:
                if shape.has_text_frame and "Sprint 2: WebGL Shaders & Sequence Diagram" in shape.text_frame.text:
                    shape.text_frame.text = "Sprint 2: WebGL Custom Shaders"
            
            # 3. Create three vertical cards and position shader texts into them
            add_card_to_back(slide, Inches(0.54), Inches(1.8), Inches(2.7), Inches(3.2))
            add_card_to_back(slide, Inches(3.65), Inches(1.8), Inches(2.7), Inches(3.2))
            add_card_to_back(slide, Inches(6.76), Inches(1.8), Inches(2.7), Inches(3.2))
            
            # Move text shapes by content
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    # Column 1
                    if "Atmospheric Glow" in text:
                        shape.left = int(Inches(0.69))
                        shape.top = int(Inches(2.0))
                        shape.width = int(Inches(2.4))
                    elif "Fresnel-effect vertex" in text:
                        shape.left = int(Inches(0.69))
                        shape.top = int(Inches(2.4))
                        shape.width = int(Inches(2.4))
                        shape.height = int(Inches(2.4))
                    # Column 2
                    elif "Neon Country Borders" in text:
                        shape.left = int(Inches(3.80))
                        shape.top = int(Inches(2.0))
                        shape.width = int(Inches(2.4))
                    elif "5 overlapping TopoJSON" in text:
                        shape.left = int(Inches(3.80))
                        shape.top = int(Inches(2.4))
                        shape.width = int(Inches(2.4))
                        shape.height = int(Inches(2.4))
                    # Column 3
                    elif "Bloom Post-Processing" in text:
                        shape.left = int(Inches(6.91))
                        shape.top = int(Inches(2.0))
                        shape.width = int(Inches(2.4))
                    elif "EffectComposer with Bloom" in text:
                        shape.left = int(Inches(6.91))
                        shape.top = int(Inches(2.4))
                        shape.width = int(Inches(2.4))
                        shape.height = int(Inches(2.4))
                        
        # Slide 20: Sprint 3 Diagrams
        elif slide_idx == 19:
            # 1. Remove picture shapes
            for shape in list(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide.shapes._spTree.remove(shape._element)
            
            # 2. Re-layout description to fill slide nicely and change title
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if "Sprint 3: Class & Sequence Diagrams" in text:
                        shape.text_frame.text = "Sprint 3: Threat Relationship Graph Model"
                    elif "The D3 Force Graph renders" in text:
                        shape.left = int(Inches(1.0))
                        shape.top = int(Inches(1.8))
                        shape.width = int(Inches(8.0))
                        shape.height = int(Inches(3.2))
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(14)
                                
        # Slide 22: Sprint 4 Performance Math & UML
        elif slide_idx == 21:
            # 1. Remove the two UML diagram pictures on the left
            for shape in list(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide.shapes._spTree.remove(shape._element)
            
            # 2. Shift all right-side elements to the center (subtract 2.04" from left coordinate of all shapes at left > 5.0)
            for shape in slide.shapes:
                if shape.left/914400 > 5.0:
                    shape.left = int(shape.left - Inches(2.04))
            
            # 3. Update Title Text
            for shape in slide.shapes:
                if shape.has_text_frame and "Sprint 4: Performance Mathematics & UML" in shape.text_frame.text:
                    shape.text_frame.text = "Sprint 4: Performance Guardrails & Event Sampling"
                    
        # Apply dark theme styling to all shapes general cases
        for shape in slide.shapes:
            # Skip slide titles and chapter headers (they are styled inside the loop, we pass is_title flag)
            is_title = (shape == title_shape) or (shape.has_text_frame and "CHAPTER" not in shape.text_frame.text and shape.top/914400 < 1.2 and shape.width/914400 > 3.0)
            style_general_shape(shape, is_title)
            
    # Save presentation
    prs.save(pptx_path)
    print("Presentation enhanced successfully.")

if __name__ == "__main__":
    pptx_path = r"c:\Users\shihe\Desktop\Final-Year-Project\predictra-threat-map\presentation\Predictra-Threat-Map.pptx"
    new_usecase_path = r"c:\Users\shihe\Desktop\Final-Year-Project\predictra-threat-map\rapport-pfe\assets\global_usecase.png"
    new_class_path = r"c:\Users\shihe\Desktop\Final-Year-Project\predictra-threat-map\rapport-pfe\assets\global_class_diagram.png"
    
    make_backup(pptx_path)
    enhance_presentation(pptx_path, new_usecase_path, new_class_path)
