import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.enum.text import PP_ALIGN

# Styling Constants
COLOR_BG = RGBColor(6, 8, 15)             # Deep Cyber Navy-Black
COLOR_CARD_BG = RGBColor(17, 22, 37)      # Sleek Slate-Blue Card Fill
COLOR_CARD_BORDER = RGBColor(30, 41, 59)  # Dark Gray Card Border
COLOR_CYAN = RGBColor(0, 240, 255)        # Neon Cyber Cyan Accent
COLOR_WHITE = RGBColor(255, 255, 255)     # Bright White for headers/highlights
COLOR_SILVER = RGBColor(203, 213, 225)    # Slate-Silver for body text
COLOR_BODY = RGBColor(168, 178, 209)      # Cool Silver-Gray for body descriptions
COLOR_HIGHLIGHT = RGBColor(56, 189, 248)  # Electric Blue Highlight

def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_transition(slide):
    # Transition XML for slow fade transition
    transition_xml = '''
    <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow">
        <p:fade />
    </p:transition>
    '''
    slide.element.insert(-1, parse_xml(transition_xml))

def add_header(slide, chapter_text, title_text):
    # Chapter Indicator
    if chapter_text:
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.73), Inches(0.3))
        tf = tx_box.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = chapter_text.upper()
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        
    # Main Slide Title
    tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.73), Inches(0.6))
    tf2 = tx_box2.text_frame
    tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.text = title_text
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    
    # Neon Divider Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.73), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_CYAN
    line.line.fill.background() # borderless

def add_card(slide, left, top, width, height, title, body_paragraphs):
    # Card container
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_CARD_BORDER
    card.line.width = Pt(1)
    
    # Text Frame
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Body
    for para_text in body_paragraphs:
        p_next = tf.add_paragraph()
        p_next.text = para_text
        p_next.font.name = "Segoe UI"
        p_next.font.size = Pt(11)
        p_next.font.color.rgb = COLOR_SILVER
        p_next.space_before = Pt(4)
    return card

def safe_add_image(slide, image_path, left, top, width, height, placeholder_name):
    if os.path.exists(image_path):
        try:
            return slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as e:
            print(f"Error loading image {image_path}: {e}")
            
    # Fallback to a styled card shape representing the image
    fallback = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fallback.fill.solid()
    fallback.fill.fore_color.rgb = COLOR_CARD_BG
    fallback.line.color.rgb = COLOR_CYAN
    fallback.line.width = Pt(1)
    tf = fallback.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[Image Asset missing: {placeholder_name}]"
    p.font.name = "Segoe UI"
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_SILVER
    p.alignment = PP_ALIGN.CENTER
    return fallback

def create_bespoke_deck(pptx_path, assets_dir, report_assets_dir):
    prs = Presentation()
    # True widescreen 16:9 dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ==================== SLIDE 1: TITLE SLIDE ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    
    # Main project title
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.33), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "PREDICTRA THREAT MAP"
    p.font.name = "Segoe UI"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Design and Development of a Real-Time Spatialized Threat Intelligence Visualization Platform"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_SILVER
    p2.alignment = PP_ALIGN.CENTER
    
    # Left Block: Presenters & Supervisor Info
    card_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.6), Inches(4.8), Inches(2.2))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = COLOR_CARD_BG
    card_left.line.color.rgb = COLOR_CARD_BORDER
    tb_left = slide.shapes.add_textbox(Inches(1.7), Inches(3.8), Inches(4.4), Inches(1.8))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0
    
    p_l1 = tf_l.paragraphs[0]
    p_l1.text = "AUTHORS"
    p_l1.font.name = "Segoe UI"
    p_l1.font.size = Pt(11)
    p_l1.font.bold = True
    p_l1.font.color.rgb = COLOR_CYAN
    
    p_l2 = tf_l.add_paragraph()
    p_l2.text = "Brahim JABALLI & Chiheb AMRI"
    p_l2.font.name = "Segoe UI"
    p_l2.font.size = Pt(14)
    p_l2.font.bold = True
    p_l2.font.color.rgb = COLOR_WHITE
    
    p_l3 = tf_l.add_paragraph()
    p_l3.text = "\nSUPERVISOR"
    p_l3.font.name = "Segoe UI"
    p_l3.font.size = Pt(11)
    p_l3.font.bold = True
    p_l3.font.color.rgb = COLOR_CYAN
    
    p_l4 = tf_l.add_paragraph()
    p_l4.text = "Mr. Anis DHAHRI — ISET Gafsa"
    p_l4.font.name = "Segoe UI"
    p_l4.font.size = Pt(13)
    p_l4.font.color.rgb = COLOR_SILVER
    
    # Right Block: Institution & Academic Info
    card_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.03), Inches(3.6), Inches(4.8), Inches(2.2))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = COLOR_CARD_BG
    card_right.line.color.rgb = COLOR_CARD_BORDER
    tb_right = slide.shapes.add_textbox(Inches(7.23), Inches(3.8), Inches(4.4), Inches(1.8))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0
    
    p_r1 = tf_r.paragraphs[0]
    p_r1.text = "HOST FIRM"
    p_r1.font.name = "Segoe UI"
    p_r1.font.size = Pt(11)
    p_r1.font.bold = True
    p_r1.font.color.rgb = COLOR_CYAN
    
    p_r2 = tf_r.add_paragraph()
    p_r2.text = "Predictra Cybersecurity"
    p_r2.font.name = "Segoe UI"
    p_r2.font.size = Pt(14)
    p_r2.font.bold = True
    p_r2.font.color.rgb = COLOR_WHITE
    
    p_r3 = tf_r.add_paragraph()
    p_r3.text = "\nACADEMIC YEAR & OPTION"
    p_r3.font.name = "Segoe UI"
    p_r3.font.size = Pt(11)
    p_r3.font.bold = True
    p_r3.font.color.rgb = COLOR_CYAN
    
    p_r4 = tf_r.add_paragraph()
    p_r4.text = "2025 / 2026  ·  Option: Computer System Development"
    p_r4.font.name = "Segoe UI"
    p_r4.font.size = Pt(12)
    p_r4.font.color.rgb = COLOR_SILVER
    
    # Jury Block
    tb_jury = slide.shapes.add_textbox(Inches(1.5), Inches(6.0), Inches(10.33), Inches(0.8))
    tf_j = tb_jury.text_frame
    p_j = tf_j.paragraphs[0]
    p_j.text = "Jury: Ms. Hadhami ISSAOUI (President)  ·  Mr. Walid HAMMEMI (Reporter)"
    p_j.font.name = "Segoe UI"
    p_j.font.size = Pt(11)
    p_j.font.italic = True
    p_j.font.color.rgb = COLOR_BODY
    p_j.alignment = PP_ALIGN.CENTER
    
    # Speaker Notes
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: Good morning, ladies and gentlemen of the jury. Welcome to our graduation defense for our "
        "Bachelor’s Degree in Information Technology. My name is Brahim Jaballi, and alongside my colleague "
        "Chiheb Amri, we are proud to present our final year project: the 'Predictra Threat Map'. This project, "
        "conducted under the supervision of Mr. Anis Dhahri at ISET Gafsa, and in collaboration with host company "
        "Predictra Cybersecurity, focuses on the design and development of a real-time spatialized threat "
        "intelligence visualization platform. We would also like to thank Ms. Hadhami Issaoui and Mr. Walid Hammemi "
        "for presiding and reporting on our defense today."
    )

    # ==================== SLIDE 2: PRESENTATION OUTLINE ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "agenda", "Presentation Outline")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(2.15), Inches(4.5), "01  Intro", 
             ["• Global landscape", "• Attack surface", "• Strategic threat intel motivation"])
    add_card(slide, Inches(3.2), Inches(2.0), Inches(2.15), Inches(4.5), "02  Pre-Study", 
             ["• Host company", "• The problematic", "• Proposed solution", "• Agile Scrum model"])
    add_card(slide, Inches(5.6), Inches(2.0), Inches(2.15), Inches(4.5), "03  Conceptual", 
             ["• System actors", "• Requirements", "• Use Case model", "• Stack comparison", "• Architecture"])
    add_card(slide, Inches(8.0), Inches(2.0), Inches(2.15), Inches(4.5), "04  Realization", 
             ["• Ingestion pipelines", "• SSE streaming math", "• WebGL Shaders", "• STIX & MITRE Map", "• Guardrails"])
    add_card(slide, Inches(10.4), Inches(2.0), Inches(2.15), Inches(4.5), "05  Conclusion", 
             ["• Achievement summary", "• Retro velocity", "• Kafka & ML scalability roadmap"])
             
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: To guide you through our presentation, we have structured our defense into five key chapters. "
        "I will begin with the Introduction and the Pre-Study, outlining the cybersecurity landscape, our problematic, "
        "and the Agile Scrum methodology. Next, we will present our Conceptual Study, including system actors, "
        "requirements, use cases, and architecture. We will then dive into the Realization phase, detailing "
        "Sprints 1 to 4 and showcasing the system interfaces. Finally, we will conclude with a retrospective, "
        "technical challenges resolved, and future perspectives."
    )

    # ==================== SLIDE 3: GENERAL INTRODUCTION ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 1 — introduction", "General Introduction")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(5.0), Inches(4.5), "The Strategic Shift to Proactive Defense", 
             ["Sun Tzu, The Art of War:",
              "\"If you know the enemy and know yourself, you need not fear the result of a hundred battles.\"",
              "\nDefenders must pivot from reactive, perimeter-based alerts (Firewalls & antivirus) to proactive external visibility. By collecting, normalizing, and spatializing threat feeds, we anticipate campaigns before they penetrate our perimeter."])
              
    add_card(slide, Inches(6.3), Inches(2.0), Inches(6.2), Inches(2.1), "1. Expanded Attack Surface", 
             ["Hyper-connectivity, cloud infrastructure, and distributed microservices have multiplied threat exposure points, making manual tracking impossible."])
    add_card(slide, Inches(6.3), Inches(4.4), Inches(6.2), Inches(2.1), "2. Evolved Adversaries", 
             ["Attackers are no longer lone hackers but structured cartels, nation-states, and automated botnets operating rapid campaigns at scale."])
             
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: Let us start with the introduction. Today, corporate networks and cloud platforms are faced "
        "with an expanded attack surface. Highly coordinated threat groups launch disruptive operations globally. "
        "As Sun Tzu wrote in The Art of War: 'If you know the enemy and know yourself, you need not fear the "
        "result of a hundred battles.' In cybersecurity, this means shifting from a reactive posture to a "
        "proactive threat intelligence model, where we understand external adversary infrastructure before a breach occurs."
    )

    # ==================== SLIDE 4: PROJECT CONTEXT & HOST ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 1 — pre-study", "Project Context & Host Organization")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(5.0), Inches(2.1), "Host: Predictra Cybersecurity", 
             ["An AI-driven Cyber Threat Intelligence (CTI) firm that focuses on democratizing raw threat datasets into immediate, visually actionable insights for enterprise SOCs."])
             
    add_card(slide, Inches(0.8), Inches(4.4), Inches(5.0), Inches(2.1), "The Dwell Time Target", 
             ["The global average dwell time for compromise detection exceeds 200 days. The primary goal of the Threat Map is to shrink this window from months to real-time seconds."])
             
    # Organization Chart Screenshot
    org_img_path = os.path.join(assets_dir, "slide_4_shape_4.png")
    safe_add_image(slide, org_img_path, Inches(6.3), Inches(2.0), Inches(6.2), Inches(4.5), "Predictra Corporate Structure")
    
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: Our host organization, Predictra Cybersecurity, is an AI-driven Cyber Threat Intelligence "
        "firm. Their goal is to democratize threat intelligence, transforming raw feeds into actionable visual "
        "insights. A critical industry metric is the compromise dwell time, which globally averages over 200 days. "
        "Security teams often detect breaches months after they occur. Our Threat Map aims to shrink this window "
        "from months to seconds by providing immediate, real-time visibility."
    )

    # ==================== SLIDE 5: THE PROBLEM ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 1 — pre-study", "The Critical Problematic")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5), "Spreadsheet Overload", 
             ["SOC analysts are overwhelmed by raw, static text logs. Staring at thousands of text rows causes fatigue and delays the detection of critical correlation anomalies."])
             
    add_card(slide, Inches(4.7), Inches(2.0), Inches(3.6), Inches(4.5), "Lack of Spatial Context", 
             ["Traditional grids omit geography. Security teams cannot visualize where attacks originate, which sectors are targeted, or the real-time vectors of global campaigns."])
             
    add_card(slide, Inches(8.6), Inches(2.0), Inches(3.9), Inches(4.5), "Data Heterogeneity", 
             ["Threat feeds are scattered across separate APIs, WebSocket channels, and raw CSV files. Formats, schemas, and protocols are incompatible, blocking immediate analysis."])
             
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: The core problem we address is threefold. First, spreadsheet overload. SOC analysts suffer "
        "from visual fatigue reading thousands of static logs. Second, a lack of spatial context. Traditional "
        "dashboards fail to show attack origins, targets, and trajectories geographically. Third, data heterogeneity. "
        "Threat feeds are scattered across multiple vendors in incompatible formats. This makes manual correlation "
        "slow, leading to delayed responses."
    )

    # ==================== SLIDE 6: PROPOSED SOLUTION ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 1 — pre-study", "The Proposed Solution: 4-Layer Architecture")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(2.7), Inches(4.5), "1. Ingestion Layer", 
             ["Asynchronous, non-blocking ingestion scraping from 9+ parallel feeds (AlienVault, URLhaus, Bitdefender, Checkpoint, etc.)."])
             
    add_card(slide, Inches(3.8), Inches(2.0), Inches(2.7), Inches(4.5), "2. Enrichment Layer", 
             ["Local geolocation dictionary mapping and 5-layer classification engine tagging targets into 9 critical industries."])
             
    add_card(slide, Inches(6.8), Inches(2.0), Inches(2.7), Inches(4.5), "3. Streaming Layer", 
             ["Server-Sent Events (SSE) broadcasting queue with a 300ms batch flush that prevents browser rendering overhead."])
             
    add_card(slide, Inches(9.8), Inches(2.0), Inches(2.7), Inches(4.5), "4. Visualization Layer", 
             ["WebGL 3D Earth rendered via Three.js at 60fps with glowing bezier arcs, D3 force graph, and MITRE ATT&CK grid."])
             
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: To solve this, we propose the Predictra Threat Map, built on a four-layer architecture. "
        "First, the Ingestion Layer harvests data asynchronously from 9+ feeds. Second, the Enrichment Layer "
        "geolocates attacks and classifies victims into 9 critical sectors. Third, the Streaming Layer broadcasts "
        "events using Server-Sent Events with a 300ms batching queue. Fourth, the Visualization Layer displays "
        "curved attack arcs on an interactive 3D WebGL Earth, alongside a STIX parser and MITRE ATT&CK grid."
    )

    # ==================== SLIDE 7: AGILE SCRUM METHODOLOGY ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 1 — pre-study", "Agile Scrum Methodology")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), "4 Sprints & Agile Rituals", 
             ["We ran a 12-week development cycle with clear team velocity metrics (25-30 SP/sprint).",
              "\nTwice-Weekly Rituals:",
              "• Thursdays: Code integrations, WebGL profiling, performance checks.",
              "• Sundays: Architectural consolidations, LaTeX report reviews.",
              "\nResult: 48+ collaborative sessions delivering a fully testable, production-ready release increment at each review."])
              
    add_card(slide, Inches(6.7), Inches(2.0), Inches(5.8), Inches(0.9), "Sprint 1: Backend", ["Ingestion engine, 9 scrapers, MongoDB schema, enrichment rules."])
    add_card(slide, Inches(6.7), Inches(3.1), Inches(5.8), Inches(0.9), "Sprint 2: Stream & 3D", ["GET /api/feed SSE endpoint, Zustand store, Three.js 3D Globe."])
    add_card(slide, Inches(6.7), Inches(4.2), Inches(5.8), Inches(0.9), "Sprint 3: STIX Workspace", ["Client-side parser, D3 relationship force graph, MITRE ATT&CK matrix."])
    add_card(slide, Inches(6.7), Inches(5.3), Inches(5.8), Inches(0.9), "Sprint 4: Analytics & Guards", ["Paginated logs, Target My IP, Excel export, Adaptive FPS sampling."])
    
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: We adopted the Agile Scrum methodology, organizing our work into 4 sprints over a 12-week period. "
        "Sprint 1 focused on backend ingestion and schemas. Sprint 2 implemented the SSE stream and the 3D globe. "
        "Sprint 3 integrated the STIX 2.1 parser and MITRE ATT&CK grid. Sprint 4 added the history browser, "
        "Target My IP, and performance guardrails. We held twice-weekly rituals on Thursdays for code merges "
        "and profiling, and Sundays for architectural consolidation and documentation."
    )

    # ==================== SLIDE 8: GLOBAL USE CASE ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 2 — conceptual study", "System Boundaries & Global Use Case")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(4.5), Inches(4.5), "Actors & Use Case Boundaries", 
             ["This model structures system actors and boundaries:",
              "\n• System Actor: External Threat Feeds",
              "Continuously pushes Indicators of Compromise (IOCs) to trigger ingestion without human intervention.",
              "\n• Human Actor: Security Analyst",
              "Interacts with the Web UI to monitor live feeds, search history, upload local STIX files, and export reports."])
              
    # Image frame
    card_img = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.8), Inches(6.0), Inches(4.9))
    card_img.fill.solid()
    card_img.fill.fore_color.rgb = COLOR_CARD_BG
    card_img.line.color.rgb = COLOR_CARD_BORDER
    
    usecase_img_path = os.path.join(report_assets_dir, "global_usecase.png")
    safe_add_image(slide, usecase_img_path, Inches(8.3), Inches(1.9), Inches(2.4), Inches(4.7), "Global Use Case")
    
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: This is our Global Use Case diagram, showing the system boundaries and interactions. The primary "
        "system actor is the set of External Threat Feeds, which automatically trigger backend ingestion. The "
        "human actor is the Security Analyst, who monitors the 3D globe and dashboard, filters threat logs, searches "
        "historical data, uploads STIX bundles, and exports reports. This diagram establishes the foundational "
        "boundaries of the Predictra Threat Map."
    )

    # ==================== SLIDE 9: SYSTEM REQUIREMENTS ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 2 — conceptual study", "System Requirements")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), "Functional Requirements (RF-1 to RF-9)", 
             ["• Ingest from 9+ feeds asynchronously.",
              "• Resolve IP geolocation and 5-layer target sector.",
              "• Stream real-time data using SSE heartbeat.",
              "• Render WebGL 3D Globe with animated attack arcs.",
              "• Client-side STIX 2.1 parsing & MITRE ATT&CK map.",
              "• Searchable log browser with pagination & filters.",
              "• \"Target My IP\" geolocation and Excel export."])
              
    add_card(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5), "Non-Functional Requirements (RNF-1 to RNF-5)", 
             ["• RNF-1 (Fluidity): Maintain 60fps Earth rendering.",
              "• RNF-2 (Memory): 10,000-event circular RingBuffer.",
              "• RNF-3 (Database): MongoDB 30-day TTL auto-expire.",
              "• RNF-4 (Limit Guard): switch to in-memory streaming and stop DB writes if MongoDB exceeds 500MB.",
              "• RNF-5 (Design): Premium dark glassmorphism interface to prevent visual analyst strain."])
              
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: We defined 9 functional requirements, denoted RF-1 to RF-9. These include multi-feed ingestion, "
        "geolocation enrichment, target sector classification, real-time SSE streaming, the 3D WebGL globe, "
        "client-side STIX parsing, MITRE ATT&CK mapping, historical log searching, and the 'Target My IP' feature. "
        "Our non-functional requirements cover 60fps fluidity, circular memory buffer, 30-day database TTL, a 500MB "
        "quota guard, and a premium dark glassmorphism layout to reduce SOC analyst fatigue."
    )

    # ==================== SLIDE 10: GLOBAL CLASS DIAGRAM ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 2 — conceptual study", "Structural Modeling: Global Class Diagram")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(4.5), Inches(4.5), "System Object Architecture", 
             ["The static structure consists of three key architectural blocks:",
              "\n1. Ingestion & Service Tier (Backend):",
              "ExpressServer, ScraperService subclasses, EnrichmentService, Mongoose models.",
              "\n2. State Store Tier (Frontend):",
              "Zustand useStreamStore holding a circular RingBuffer and perfTelemetry loop.",
              "\n3. React Views Layer:",
              "Dashboard, STIXWorkspace, AnalyticsPage subscribing to store slices."])
              
    # Image frame
    card_img2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.6), Inches(1.8), Inches(6.9), Inches(4.5))
    card_img2.fill.solid()
    card_img2.fill.fore_color.rgb = COLOR_CARD_BG
    card_img2.line.color.rgb = COLOR_CARD_BORDER
    
    class_img_path = os.path.join(report_assets_dir, "global_class_diagram.png")
    safe_add_image(slide, class_img_path, Inches(5.8), Inches(2.0), Inches(6.5), Inches(4.1), "Global Class Diagram")
    
    slide.notes_slide.notes_text_frame.text = (
        "Chiheb: This is our Global Class Diagram, showing the structural design of the platform. In the Backend "
        "tier, the Express Server schedules multiple Scraper Services, which feed into the Enrichment Service to "
        "resolve geolocations and sectors. The Zustand Store manages frontend state with a circular RingBuffer and "
        "performance telemetry. The React Views, such as the Dashboard, STIX Workspace, and Analytics page, "
        "subscribe directly to Zustand slices for reactivity."
    )

    # ==================== SLIDE 11: TECHNOLOGY STACK ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 2 — conceptual study", "Technology Stack Comparative Analysis")
    
    # Table headers
    add_card(slide, Inches(0.8), Inches(2.0), Inches(2.7), Inches(0.8), "Layer", [])
    add_card(slide, Inches(3.6), Inches(2.0), Inches(2.7), Inches(0.8), "Selected Stack", [])
    add_card(slide, Inches(6.4), Inches(2.0), Inches(2.7), Inches(0.8), "Alternative", [])
    add_card(slide, Inches(9.2), Inches(2.0), Inches(3.3), Inches(0.8), "Architectural Rationale", [])
    
    # Row 1
    add_card(slide, Inches(0.8), Inches(2.9), Inches(2.7), Inches(0.8), "Backend Runtime", [])
    add_card(slide, Inches(3.6), Inches(2.9), Inches(2.7), Inches(0.8), "Node.js / Express 5", [])
    add_card(slide, Inches(6.4), Inches(2.9), Inches(2.7), Inches(0.8), "Java Spring Boot", [])
    add_card(slide, Inches(9.2), Inches(2.9), Inches(3.3), Inches(0.8), "Asynchronous loop; high SSE concurrency", [])
    
    # Row 2
    add_card(slide, Inches(0.8), Inches(3.8), Inches(2.7), Inches(0.8), "State Management", [])
    add_card(slide, Inches(3.6), Inches(3.8), Inches(2.7), Inches(0.8), "Zustand 5", [])
    add_card(slide, Inches(6.4), Inches(3.8), Inches(2.7), Inches(0.8), "Redux Toolkit", [])
    add_card(slide, Inches(9.2), Inches(3.8), Inches(3.3), Inches(0.8), "Reactive store outside React render loop", [])
    
    # Row 3
    add_card(slide, Inches(0.8), Inches(4.7), Inches(2.7), Inches(0.8), "3D Graphics engine", [])
    add_card(slide, Inches(3.6), Inches(4.7), Inches(2.7), Inches(0.8), "Three.js / React Three Fiber", [])
    add_card(slide, Inches(6.4), Inches(4.7), Inches(2.7), Inches(0.8), "Vanilla HTML Canvas", [])
    add_card(slide, Inches(9.2), Inches(4.7), Inches(3.3), Inches(0.8), "GPU acceleration, custom GLSL shaders", [])
    
    # Row 4
    add_card(slide, Inches(0.8), Inches(5.6), Inches(2.7), Inches(0.8), "Database Tier", [])
    add_card(slide, Inches(3.6), Inches(5.6), Inches(2.7), Inches(0.8), "MongoDB Atlas (BSON)", [])
    add_card(slide, Inches(6.4), Inches(5.6), Inches(2.7), Inches(0.8), "PostgreSQL", [])
    add_card(slide, Inches(9.2), Inches(5.6), Inches(3.3), Inches(0.8), "Flexible schemaless BSON for raw feeds", [])
    
    slide.notes_slide.notes_text_frame.text = (
        "Chiheb: For the backend runtime, we selected Node.js and Express 5 over Java Spring Boot due to its "
        "non-blocking asynchronous event loop, which is perfect for streaming SSE connections under low memory footprints. "
        "For state management, we chose Zustand 5 over Redux for high-speed reactive updates outside the React DOM tree, "
        "preserving WebGL rendering cycles. For graphics, we used Three.js and React Three Fiber to build custom shaders with "
        "GPU acceleration. For the database, we chose MongoDB Atlas for its schemaless structure."
    )

    # ==================== SLIDE 12: SYSTEM ARCHITECTURE ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 2 — conceptual study", "System Architecture")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5), "Frontend (Vercel CDN)", 
             ["• React Single Page App served via global Edge networks.",
              "\n• Zustand state-store for real-time reactivity.",
              "\n• Vercel configuration proxy routing rules to prevent CORS blocks."])
              
    add_card(slide, Inches(4.7), Inches(2.0), Inches(3.6), Inches(4.5), "Backend (Cloud VM)", 
             ["• Express REST APIs + SSE event broadcast loop.",
              "\n• PM2 Clustering to distribute scraping workload across multiple CPU cores.",
              "\n• Ingestion schedulers managing 9 parallel scrapers."])
              
    add_card(slide, Inches(8.6), Inches(2.0), Inches(3.9), Inches(4.5), "Database (MongoDB)", 
             ["• Cloud-hosted replicated BSON cluster.",
              "\n• Automated 30-day Time-To-Live (TTL) storage pruning indexes.",
              "\n• Storage monitoring: switches to in-memory streaming if DB size hits 500MB."])
              
    slide.notes_slide.notes_text_frame.text = (
        "Chiheb: Our system architecture is split into three main tiers. The Frontend is built in React and "
        "served via Vercel CDN, utilizing API routing proxies to eliminate cross-origin resource sharing blocks. "
        "The Backend runs on a Cloud VM using PM2 process clustering across CPU cores for automatic crash recovery. "
        "The Database tier uses MongoDB Atlas with a replicated cluster, automatic TTL pruning, and size-based quota guards."
    )

    # ==================== SLIDE 13: SPRINT 1 & 2 REALIZATION ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 3 — sprints 1 & 2", "Sprint 1 & 2: Ingestion & 3D Math")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5), "Ingestion Scrapers (9 Feeds)", 
             ["• Feeds: Checkpoint SSE, Bitdefender WebSockets, FortiGuard REST, URLhaus, Kaspersky, OTX, RansomWatch, C2 Tracker, MISP Galaxy.",
              "\n• Compact BSON schema compresses JSON models, shrinking database footprint by 60% (200B → 80B per threat document)."])
              
    add_card(slide, Inches(4.7), Inches(2.0), Inches(3.6), Inches(4.5), "5-Layer Enrichment logic", 
             ["• Layered classification rules (port map, keywords, regex tags, MISP actors, feed defaults).",
              "\n• Asynchronous RDAP owner lookups with a 3-second timeout cache to prevent API blocks."])
              
    add_card(slide, Inches(8.6), Inches(2.0), Inches(3.9), Inches(4.5), "3D Spherical Coordinate Math", 
             ["• Convert coordinates (lat/lon) to 3D vectors via spherical trig.",
              "\n• Attack arcs compute using Spherical Linear Interpolation (SLERP), combined with quadratic Bezier curves lifting longer paths higher above Earth."])
              
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: In Sprints 1 and 2, we built the core ingestion and 3D visualization. We integrated 9 distinct "
        "feeds and compressed the BSON schema from 200 bytes to 80 bytes. The enrichment logic resolves geocoding "
        "and victim sector using a 5-layer engine. For the 3D globe, geographic coordinates are converted to Cartesian "
        "vectors, and attack arcs are animated using SLERP and Bezier elevation math."
    )

    # ==================== SLIDE 14: SPRINT 2 CUSTOM SHADERS ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 3 — sprint 2", "Sprint 2: WebGL Custom Shaders")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5), "1. Atmospheric Glow", 
             ["• Custom Fresnel-effect vertex and fragment shaders.",
              "\n• Dot-product vector math computes viewing angle falloff, rendering a neon cyberblue glow around the Earth's silhouette."])
              
    add_card(slide, Inches(4.7), Inches(2.0), Inches(3.6), Inches(4.5), "2. Neon Country Borders", 
             ["• 5 overlapping TopoJSON vector layers placed at microscopic radial scales (R = 1.002 to 1.010).",
              "\n• Decreasing opacity scales simulate a glowing country border outline effect on the globe mesh."])
              
    add_card(slide, Inches(8.6), Inches(2.0), Inches(3.9), Inches(4.5), "3. Bloom Post-Processing", 
             ["• Three.js EffectComposer with UnrealBloomPass shader.",
              "\n• Luminance threshold set to 0.15 makes the severity-colored attack arcs glow like fiber-optic light strings."])
              
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: To make the globe visually stunning, we wrote custom WebGL shaders. The Atmospheric Glow "
        "uses a Fresnel-effect vertex and fragment shader. Neon Country Borders are rendered using 5 overlapping "
        "TopoJSON vector layers with decreasing opacities. Finally, Bloom Post-Processing applies a luminosity "
        "threshold to make severity-colored attack arcs glow like neon light fibers."
    )

    # ==================== SLIDE 15: SPRINT 3 & 4 REALIZATION ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 4 — sprints 3 & 4", "Sprint 3 & 4: STIX, MITRE & Performance")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5), "STIX 2.1 Workspace", 
             ["• Client-side JSON parser processes files in browser (zero server upload ensures analyst privacy).",
              "\n• Interactive D3 Force-Directed Graph shows relationships (Verlet integration)."])
              
    add_card(slide, Inches(4.7), Inches(2.0), Inches(3.6), Inches(4.5), "MITRE ATT&CK Matrix", 
             ["• Maps parsed STIX indicators to 12 ATT&CK kill chain phases.",
              "\n• Real-time aggregations highlight critical adversary tactics."])
              
    add_card(slide, Inches(8.6), Inches(2.0), Inches(3.9), Inches(4.5), "Performance Guardrails", 
             ["• Circular RingBuffer manages memory allocation in O(1) time.",
              "\n• Adaptive Event Sampling checks FPS. If FPS drops below 55, it drops 50% of visual arcs; if below 30, it drops 75% to preserve smoothness."])
              
    slide.notes_slide.notes_text_frame.text = (
        "Chiheb: Sprints 3 and 4 introduced STIX parsing, MITRE mapping, and performance guards. The STIX 2.1 "
        "parser runs completely in the browser for privacy, mapping techniques to 12 MITRE phases. The D3 graph "
        "uses Verlet integration to layout nodes. To maintain 60fps, we built a circular RingBuffer and an "
        "Adaptive Event Sampling loop that drops up to 75% of visual events under heavy CPU load."
    )

    # ==================== SLIDE 16: INTERFACE OVERVIEW ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 4 — realization", "Application Interfaces Overview")
    
    # 2x2 collage of dashboard screenshots
    # Card 1: 3D Live Threat Map
    card1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.2))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_CARD_BG
    card1.line.color.rgb = COLOR_CARD_BORDER
    tb1 = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.2), Inches(1.8))
    tb1.text_frame.word_wrap = True
    p_tb1 = tb1.text_frame.paragraphs[0]
    p_tb1.text = "1. 3D Live Threat Map"
    p_tb1.font.name = "Segoe UI"
    p_tb1.font.size = Pt(14)
    p_tb1.font.bold = True
    p_tb1.font.color.rgb = COLOR_CYAN
    p_tb1_b = tb1.text_frame.add_paragraph()
    p_tb1_b.text = "\nInteractive WebGL Earth showing real-time arcs, glowing outlines, and country highlight panels."
    p_tb1_b.font.name = "Segoe UI"
    p_tb1_b.font.size = Pt(11)
    p_tb1_b.font.color.rgb = COLOR_BODY
    
    img1_path = os.path.join(assets_dir, "slide_23_shape_2.png")
    safe_add_image(slide, img1_path, Inches(4.3), Inches(1.9), Inches(2.0), Inches(2.0), "3D Globe Screen")
    
    # Card 2: System Dashboard
    card2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.93), Inches(1.8), Inches(5.6), Inches(2.2))
    card2.fill.solid()
    card2.fill.fore_color.rgb = COLOR_CARD_BG
    card2.line.color.rgb = COLOR_CARD_BORDER
    tb2 = slide.shapes.add_textbox(Inches(7.13), Inches(2.0), Inches(3.2), Inches(1.8))
    tb2.text_frame.word_wrap = True
    p_tb2 = tb2.text_frame.paragraphs[0]
    p_tb2.text = "2. System Dashboard"
    p_tb2.font.name = "Segoe UI"
    p_tb2.font.size = Pt(14)
    p_tb2.font.bold = True
    p_tb2.font.color.rgb = COLOR_CYAN
    p_tb2_b = tb2.text_frame.add_paragraph()
    p_tb2_b.text = "\nTactical command center displaying volume sparklines, 2D flat map toggles, and country statistics."
    p_tb2_b.font.name = "Segoe UI"
    p_tb2_b.font.size = Pt(11)
    p_tb2_b.font.color.rgb = COLOR_BODY
    
    img2_path = os.path.join(assets_dir, "slide_23_shape_5.png")
    safe_add_image(slide, img2_path, Inches(10.43), Inches(1.9), Inches(2.0), Inches(2.0), "System Dashboard Screen")
    
    # Card 3: STIX Workspace
    card3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.2))
    card3.fill.solid()
    card3.fill.fore_color.rgb = COLOR_CARD_BG
    card3.line.color.rgb = COLOR_CARD_BORDER
    tb3 = slide.shapes.add_textbox(Inches(1.0), Inches(4.6), Inches(3.2), Inches(1.8))
    tb3.text_frame.word_wrap = True
    p_tb3 = tb3.text_frame.paragraphs[0]
    p_tb3.text = "3. STIX Workspace"
    p_tb3.font.name = "Segoe UI"
    p_tb3.font.size = Pt(14)
    p_tb3.font.bold = True
    p_tb3.font.color.rgb = COLOR_CYAN
    p_tb3_b = tb3.text_frame.add_paragraph()
    p_tb3_b.text = "\nJSON bundle ingestion, relationship force graphs, and technique mapping grids."
    p_tb3_b.font.name = "Segoe UI"
    p_tb3_b.font.size = Pt(11)
    p_tb3_b.font.color.rgb = COLOR_BODY
    
    img3_path = os.path.join(assets_dir, "slide_23_shape_8.png")
    safe_add_image(slide, img3_path, Inches(4.3), Inches(4.5), Inches(2.0), Inches(2.0), "STIX Workspace Screen")
    
    # Card 4: Analytics Dashboard
    card4 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.93), Inches(4.4), Inches(5.6), Inches(2.2))
    card4.fill.solid()
    card4.fill.fore_color.rgb = COLOR_CARD_BG
    card4.line.color.rgb = COLOR_CARD_BORDER
    tb4 = slide.shapes.add_textbox(Inches(7.13), Inches(4.6), Inches(3.2), Inches(1.8))
    tb4.text_frame.word_wrap = True
    p_tb4 = tb4.text_frame.paragraphs[0]
    p_tb4.text = "4. Analytics Dashboard"
    p_tb4.font.name = "Segoe UI"
    p_tb4.font.size = Pt(14)
    p_tb4.font.bold = True
    p_tb4.font.color.rgb = COLOR_CYAN
    p_tb4_b = tb4.text_frame.add_paragraph()
    p_tb4_b.text = "\nAggregation analytics showing targeted sector metrics, actor trends, and country matrices."
    p_tb4_b.font.name = "Segoe UI"
    p_tb4_b.font.size = Pt(11)
    p_tb4_b.font.color.rgb = COLOR_BODY
    
    img4_path = os.path.join(assets_dir, "slide_23_shape_11.png")
    safe_add_image(slide, img4_path, Inches(10.43), Inches(4.5), Inches(2.0), Inches(2.0), "Analytics Dashboard Screen")
    
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: Here is an overview of our application interfaces. The 3D Live Threat Map displays the "
        "rotating Earth with glowing arcs. The System Dashboard provides real-time metrics, country statistics, "
        "and a 2D flat map. The STIX Intelligence Workspace features the D3 graph and the MITRE ATT&CK grid. "
        "Finally, the Analytics Dashboard aggregates threat metrics, sector distributions, and trends."
    )

    # ==================== SLIDE 17: TECHNICAL CHALLENGES ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 4 — realization", "Technical Challenges Resolved")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5), "1. Visual Interface Lag", 
             ["Cause: Browser canvas bottleneck when processing >500 attacks per second.",
              "\nFix: Implemented a 300ms event grouping queue, capped active arcs to 150, and created our custom adaptive FPS sampling engine to drop visual arcs under heavy CPU load."])
              
    add_card(slide, Inches(4.7), Inches(2.0), Inches(3.6), Inches(4.5), "2. Database Limits", 
             ["Cause: MongoDB Atlas free tier quotas are easily exhausted by continuous threat feeds.",
              "\nFix: Field name compression (BSON models), automated 30-day index TTL pruning, and a 500MB quota guard that stops DB writes to stream data in-memory only."])
              
    add_card(slide, Inches(8.6), Inches(2.0), Inches(3.9), Inches(4.5), "3. CORS Deployment Blocks", 
             ["Cause: Browser security prevents React on localhost/Vercel from reading the API on a separate cloud VM.",
              "\nFix: Implemented Vercel reverse-proxy redirect rules (`vercel.json`), routing all backend API calls through internal paths to bypass CORS headers."])
              
    slide.notes_slide.notes_text_frame.text = (
        "Brahim: During the realization, we successfully resolved three major technical hurdles. First, visual lag "
        "was mitigated by event batching, arc limits, and adaptive sampling. Second, database capacity constraints "
        "were addressed using BSON compression, MongoDB TTLs, and a 500MB write guard. Third, CORS deployment blocks "
        "were resolved via Vercel proxy configuration rules."
    )

    # ==================== SLIDE 18: FUTURE ROADMAP ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 5 — conclusion & perspectives", "Future Perspectives & Roadmap")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5), "Kafka Distributed Ingestion", 
             ["Decouple scrapers from the server. Scrapers publish events directly to Apache Kafka topics, and the Express runtime consumes streams from Kafka. This prevents server crashes during high traffic."])
              
    add_card(slide, Inches(4.7), Inches(2.0), Inches(3.6), Inches(4.5), "Predictive ML Extensions", 
             ["Train ML models (LSTM/time-series) on historical threat databases to forecast attack spikes and sector-targeted waves. These predictions can be overlayed as visual heatmaps on the 3D globe."])
              
    add_card(slide, Inches(8.6), Inches(2.0), Inches(3.9), Inches(4.5), "Sharded Database Clustering", 
             ["Partition MongoDB clusters geographically by country code. Queries for regional metrics will route only to the corresponding server slice, minimizing indexing latency and scaling indefinitely."])
              
    slide.notes_slide.notes_text_frame.text = (
        "Chiheb: Looking ahead, we have designed a future roadmap focusing on three phases. First, Kafka Distributed Ingestion "
        "to decouple backend scraping from the client streaming runtime. Second, Predictive ML Extensions to forecast attack "
        "waves and overlay predictive heatmaps on the Earth. Third, Sharded Geo-Database Clustering to partition MongoDB "
        "by country code to speed up regional queries."
    )

    # ==================== SLIDE 19: GENERAL CONCLUSION ====================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_transition(slide)
    add_header(slide, "chapter 5 — conclusion & perspectives", "General Conclusion")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.2), "Technical & Operational Impact", 
             ["• Ingestion & Streaming:",
              "Successful integration of 9+ live threat APIs, normalized schema, and high-performance Server-Sent Events stream.",
              "\n• 3D WebGL UI Visuals:",
              "A premium 60fps interactive Earth showing curved attack vectors, neon outlines, and custom GLSL atmospheric glows.",
              "\n• SOC Analyst Efficiency:",
              "Reduces compromise dwell time from 200+ days to seconds by replacing spreadsheet lists with high-impact visual spatial coordinates."])
              
    add_card(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.2), "Academic & Professional Growth", 
             ["• Agile Scrum Application:",
              "Twice-weekly rituals, code merges, story points track, and demonstrable increments delivered at each review.",
              "\n• Systems Engineering:",
              "Applied rigorous architectural structures, state store abstractions, and database optimization strategies.",
              "\n• Cybersecurity Focus:",
              "Deep study of Indicator of Compromise (IOC) types, STIX 2.1 schemas, and MITRE ATT&CK kill chain matrices."])
              
    # Bottom Banner: Thank You
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.4), Inches(11.73), Inches(0.6))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_CARD_BG
    banner.line.color.rgb = COLOR_CYAN
    banner.line.width = Pt(1)
    tf_b = banner.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = "Thank You for Your Attention  —  Questions Welcome"
    p_b.font.name = "Segoe UI"
    p_b.font.size = Pt(14)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_WHITE
    p_b.alignment = PP_ALIGN.CENTER
    
    slide.notes_slide.notes_text_frame.text = (
        "Chiheb: In conclusion, the Predictra Threat Map successfully bridges the visibility gap in cyber "
        "defense. It integrates 9+ feeds, feeds a 3D WebGL globe at 60fps, and features a STIX 2.1 parser. It "
        "reduces dwell time from months to seconds. This project demonstrates a rigorous academic and engineering "
        "contribution. Thank you for your time, and we welcome any questions you may have."
    )

    # Save presentation
    prs.save(pptx_path)
    print("New bespoke presentation generated successfully.")

if __name__ == "__main__":
    pptx_path = r"c:\Users\shihe\Desktop\Final-Year-Project\predictra-threat-map\presentation\Predictra-Threat-Map-Bespoke.pptx"
    assets_dir = r"c:\Users\shihe\Desktop\Final-Year-Project\predictra-threat-map\presentation\assets"
    report_assets_dir = r"c:\Users\shihe\Desktop\Final-Year-Project\predictra-threat-map\rapport-pfe\assets"
    
    create_bespoke_deck(pptx_path, assets_dir, report_assets_dir)
